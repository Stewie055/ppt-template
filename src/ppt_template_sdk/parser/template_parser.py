from __future__ import annotations

import re
from dataclasses import dataclass, field
from typing import Iterable

from pptx.enum.shapes import MSO_SHAPE_TYPE

from ..models.placeholder import Placeholder


PLACEHOLDER_RE = re.compile(r"^ph:(text|image|table|chart):([\w\.-]+)$")


@dataclass(slots=True)
class ParsedTemplate:
    placeholders: list[Placeholder] = field(default_factory=list)
    invalid_placeholders: list[str] = field(default_factory=list)
    text_field_paths: set[str] = field(default_factory=set)


def iter_shapes(shape_collection) -> Iterable:
    for shape in shape_collection:
        yield shape
        if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
            yield from iter_shapes(shape.shapes)


def parse_presentation(presentation, text_field_pattern: str) -> ParsedTemplate:
    parsed = ParsedTemplate()
    field_re = re.compile(text_field_pattern)
    for slide_index, slide in enumerate(presentation.slides):
        for shape in iter_shapes(slide.shapes):
            name = getattr(shape, "name", "") or ""
            if name.startswith("ph:"):
                match = PLACEHOLDER_RE.match(name)
                if not match:
                    parsed.invalid_placeholders.append(
                        f"slide {slide_index + 1} shape {getattr(shape, 'shape_id', '?')}: invalid placeholder name '{name}'"
                    )
                else:
                    parsed.placeholders.append(
                        Placeholder(
                            type=match.group(1),
                            key=match.group(2),
                            slide_index=slide_index,
                            shape_id=getattr(shape, "shape_id", -1),
                            shape_name=name,
                            left=getattr(shape, "left", 0),
                            top=getattr(shape, "top", 0),
                            width=getattr(shape, "width", 0),
                            height=getattr(shape, "height", 0),
                            shape=shape,
                        )
                    )

            if getattr(shape, "has_text_frame", False):
                parsed.text_field_paths.update(field_re.findall(shape.text_frame.text))
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    for cell in row.cells:
                        parsed.text_field_paths.update(field_re.findall(cell.text))
    return parsed

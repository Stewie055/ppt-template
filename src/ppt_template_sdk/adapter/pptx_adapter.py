from __future__ import annotations

from typing import Optional

from pptx import Presentation

from ..exceptions import ShapeOperationError
from ..models.content import ChartContent, ImageContent, TableContent, TextContent
from ..parser.template_parser import iter_shapes


class PptxAdapter:
    def load(self, template_path: Optional[str] = None, template_bytes: Optional[bytes] = None) -> Presentation:
        if bool(template_path) == bool(template_bytes):
            raise ValueError("exactly one of template_path or template_bytes must be provided")
        if template_path:
            return Presentation(template_path)
        from io import BytesIO

        return Presentation(BytesIO(template_bytes))

    @staticmethod
    def save_to_bytes(presentation) -> bytes:
        from io import BytesIO

        buffer = BytesIO()
        presentation.save(buffer)
        return buffer.getvalue()

    @staticmethod
    def save_to_path(presentation, output_path: str) -> None:
        presentation.save(output_path)

    @staticmethod
    def write_content(presentation, placeholder, content) -> None:
        slide = presentation.slides[placeholder.slide_index]
        shape = placeholder.shape
        if isinstance(content, TextContent):
            if not getattr(shape, "has_text_frame", False):
                raise ShapeOperationError(f"shape '{placeholder.shape_name}' cannot accept text content")
            shape.text = content.text
            return
        if isinstance(content, (ImageContent, ChartContent)):
            slide.shapes.add_picture(content.image_path, placeholder.left, placeholder.top, placeholder.width, placeholder.height)
            PptxAdapter._remove_shape(shape)
            return
        if isinstance(content, TableContent):
            if getattr(shape, "has_table", False):
                if PptxAdapter._rewrite_table(shape.table, content):
                    return
            rows, cols, grid = PptxAdapter._build_table_grid(content)
            table_shape = slide.shapes.add_table(rows, cols, placeholder.left, placeholder.top, placeholder.width, placeholder.height)
            table = table_shape.table
            for row_index, row_values in enumerate(grid):
                for col_index, value in enumerate(row_values):
                    table.cell(row_index, col_index).text = value
            PptxAdapter._remove_shape(shape)
            return
        raise ShapeOperationError(f"unsupported content type '{type(content).__name__}'")

    @staticmethod
    def replace_text_fields(presentation, rendered_shape_ids: set[int], resolver, pattern: str) -> list[str]:
        import re

        field_re = re.compile(pattern)
        warnings: list[str] = []
        for slide in presentation.slides:
            for shape in iter_shapes(slide.shapes):
                if getattr(shape, "shape_id", None) in rendered_shape_ids:
                    continue
                if getattr(shape, "has_text_frame", False):
                    shape.text = PptxAdapter._replace_text(shape.text_frame.text, field_re, resolver, warnings)
                if getattr(shape, "has_table", False):
                    for row in shape.table.rows:
                        for cell in row.cells:
                            cell.text = PptxAdapter._replace_text(cell.text, field_re, resolver, warnings)
        return warnings

    @staticmethod
    def _replace_text(text: str, field_re, resolver, warnings: list[str]) -> str:
        def repl(match):
            path = match.group(1)
            value = resolver(path)
            if value is None:
                warnings.append(f"missing text field '{path}'")
                return ""
            return str(value)

        return field_re.sub(repl, text)

    @staticmethod
    def _build_table_grid(content: TableContent) -> tuple[int, int, list[list[str]]]:
        grid: list[list[str]] = []
        if content.headers:
            grid.append([str(value) for value in content.headers])
        for row in content.rows:
            grid.append([str(value) for value in row])
        cols = max((len(row) for row in grid), default=1)
        normalized = [row + [""] * (cols - len(row)) for row in grid] or [[""]]
        return len(normalized), cols, normalized

    @staticmethod
    def _rewrite_table(table, content: TableContent) -> bool:
        rows, cols, grid = PptxAdapter._build_table_grid(content)
        if len(table.rows) != rows or len(table.columns) != cols:
            return False
        for row_index, row_values in enumerate(grid):
            for col_index, value in enumerate(row_values):
                table.cell(row_index, col_index).text = value
        return True

    @staticmethod
    def _remove_shape(shape) -> None:
        parent = shape.element.getparent()
        if parent is None:
            raise ShapeOperationError("unable to remove placeholder shape after overlay insertion")
        parent.remove(shape.element)

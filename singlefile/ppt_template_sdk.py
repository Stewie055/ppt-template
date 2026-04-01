"""单文件版 `ppt_template_sdk`。

这个文件是 PPT 模板渲染 SDK 的完整单文件发行版，适合直接复制到业务项目中维护。
如果你是第一次接触这个库，先记住这 5 个概念：

1. template:
    一个普通 `.pptx` 文件。
2. placeholder:
    你在模板里选中某个 shape，把它的 `shape.name` 命名成 `ph:<type>:<key>`，
    这个 shape 就成为一个可渲染占位块。
3. type:
    placeholder 的内容类型，决定 renderer 应该返回什么内容对象。
4. key:
    placeholder 的业务名字，用来在 `RendererRegistry` 里找到对应 renderer。
5. renderer:
    一段 Python 函数或类，它接收 `placeholder + context`，返回一个 `Content` 对象。

模板协议：
    shape.name 必须符合 `ph:<type>:<key>`，例如：

    - `ph:text:title`
    - `ph:image:logo`
    - `ph:table:risk_table`
    - `ph:chart:sales_chart`

`type` 与 renderer 返回值的对应关系：
    - `text`
        用于文本类 placeholder。renderer 返回 `TextContent`。
        SDK 会把文本写回原文本框，并尽量保留原 placeholder 的主样式。
    - `image`
        用于图片区域 placeholder。renderer 返回 `ImageContent`。
        SDK 会在该区域插入图片；这类内容通常不是直接修改原 shape。
    - `chart`
        当前按图片方式写回。renderer 返回 `ChartContent`。
        它和 `image` 类似，本质上是“提供一张图，写回到指定区域”。
    - `table`
        用于表格区域 placeholder。renderer 有两种返回方式：
        1. `TableContent`：整表替换，适合你一次生成整张表的数据。
        2. `TableCellsContent`：局部更新 cell，适合只改原生表格中的少量单元格。
        需要覆盖字体样式时，使用 `cell(...)` helper。

为什么 renderer 返回 `Content`，而不是直接改 shape：
    - 业务层只需要描述“要放什么内容”，不用关心 `python-pptx` 的底层写回细节。
    - 不同 `type` 的写回方式不同，尤其是 `image` / `chart` 往往不是原位修改。
    - SDK 可以统一做类型校验、样式保持和错误处理。

接口列表：
    PptTemplateEngine:
        模板渲染主入口。负责加载模板、解析 placeholder、调用 renderer、执行文本替换并输出 PPT。
    RendererRegistry:
        管理 `key -> renderer` 的映射。注意它不是按 `type` 注册，而是按 `key` 注册。
    BaseRenderer:
        类式 renderer 的基类。业务方可以继承它并实现 `render()`。
    RenderContext:
        统一的业务上下文对象。用于字段替换和 renderer 读取业务数据。
    TextReplacer:
        独立文本替换模块。用于已有 PPT 的普通文本框和表格单元格字段替换。
    PptOperations:
        PPT 结构操作模块。用于删页、插页、section 管理、表格结构调整和局部 cell 更新。
    EngineOptions:
        渲染引擎配置对象。
    TextContent / ImageContent / TableContent / TableCellsContent / ChartContent:
        renderer 返回的标准内容类型。
    cell:
        表格 cell helper。用于同时提供文本和值级字体样式覆盖。
    RenderResult / ValidationReport / TextReplaceResult:
        渲染、校验和文本替换的结构化结果。
    PptTemplateSdkError 及其子类:
        SDK 的统一异常体系。

最小流程：
    1. 在 PPT 模板中给某个 shape 命名 `ph:<type>:<key>`。
    2. 在 Python 中用 `RendererRegistry` 为这个 `key` 注册 renderer。
    3. renderer 根据业务数据返回对应的 `Content`。
    4. 调用 `PptTemplateEngine.render()` 输出渲染后的 PPT。

Example:
    ```python
    from ppt_template_sdk import (
        PptTemplateEngine,
        RenderContext,
        RendererRegistry,
        TableCellsContent,
        TextContent,
        cell,
    )

    registry = RendererRegistry()

    @registry.renderer("title")
    def render_title(placeholder, context):
        return TextContent(text=context.get_value("project.name", "未命名项目"))

    @registry.renderer("risk_table")
    def render_risk_table(placeholder, context):
        return TableCellsContent(cells={(1, 0): "现金流", (1, 1): cell("高", color="FF0000")})

    engine = PptTemplateEngine(registry=registry)
    result = engine.render(
        template_path="report_template.pptx",
        output_path="report_output.pptx",
        context=RenderContext(data={"project": {"name": "北极星"}}),
    )
    ```

Notes:
    - 本文件覆盖当前单文件版的全部公开能力。
    - 适合复制一个文件进项目的使用方式。
    - 不建议与包版在同一解释器环境中混用。
"""

from __future__ import annotations

import re
import uuid
from collections import Counter, defaultdict
from dataclasses import dataclass, field
from typing import Any, Callable, Iterable, Optional, Union

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt


_MISSING = object()
P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
P14_NS = "http://schemas.microsoft.com/office/powerpoint/2010/main"
R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
SECTION_EXT_URI = "{521415D9-36F7-43E2-AB2F-B90AF26B5E84}"
DEFAULT_SECTION_NAME = "Section 1"
PLACEHOLDER_RE = re.compile(r"^ph:(text|image|table|chart):([\w\.-]+)$")


class PptTemplateSdkError(Exception):
    """SDK 所有自定义异常的基类。

    当调用方不关心具体错误细分时，可以统一捕获这个异常。
    """


class TemplateParseError(PptTemplateSdkError):
    """模板解析阶段发生错误时抛出。

    这类错误通常说明模板结构本身存在问题，而不是业务 renderer 的问题。
    """


class PlaceholderFormatError(TemplateParseError):
    """占位块命名不符合 ``ph:<type>:<key>`` 规范时抛出。

    适合提示模板作者修正 `shape.name`。
    """


class DuplicatePlaceholderError(PptTemplateSdkError):
    """检测到重复占位块 key 且当前策略不允许时抛出。

    常见于 ``duplicate_key_policy="error"`` 的场景。
    """


class RendererNotFoundError(PptTemplateSdkError):
    """模板中的占位块缺少对应 renderer 时抛出。

    这通常意味着模板 key 和业务侧注册 key 不一致。
    """


class ContentTypeMismatchError(PptTemplateSdkError):
    """renderer 返回的内容类型与占位块类型不匹配时抛出。

    例如模板是 ``ph:image:cover``，renderer 却返回 ``TextContent``。
    """


class ShapeOperationError(PptTemplateSdkError):
    """对底层 shape 执行写回或定位操作失败时抛出。

    常见于目标 shape 不支持某种写回方式，或按给定定位信息找不到 shape。
    """


class FieldReplaceError(PptTemplateSdkError):
    """文本字段替换阶段发生异常时抛出。

    当字段替换遍历、取值或写回过程失败时，调用方可捕获此异常。
    """


class OperationError(PptTemplateSdkError):
    """执行 slide、section 或表格结构操作失败时抛出。

    例如索引越界、非法 merge 区域，或对已合并表格删行删列。
    """


@dataclass
class RenderContext:
    """承载模板渲染和文本替换所需的业务上下文。

    `data` 主要用于模板字段替换和轻量 renderer 取值，`extras` 适合挂复杂对象、
    聚合对象或服务实例。

    Args:
        data: 主数据源，可为 `dict`、普通对象、dataclass 或 pydantic 风格对象。
        extras: 供 renderer 直接读取的附加对象映射。

    Example:
        ```python
        context = RenderContext(
            data={"project": {"name": "北极星"}},
            extras={"report": report},
        )
        ```
    """

    data: Any
    extras: dict[str, Any] = field(default_factory=dict)

    def get_value(self, path: str, default: Any = None) -> Any:
        """按点路径读取 `data` 中的值。

        Args:
            path: 点路径，例如 `"project.name"` 或 `"items.0.title"`。
            default: 路径不存在时返回的默认值。

        Returns:
            解析后的值；若路径不存在，则返回 `default`。

        Example:
            ```python
            value = context.get_value("project.name", "未命名项目")
            ```
        """

        current = self.data
        for part in path.split("."):
            current = self._resolve_part(current, part, _MISSING)
            if current is _MISSING:
                return default
        return current

    def has_value(self, path: str) -> bool:
        """判断给定点路径是否存在有效值。

        Args:
            path: 需要检查的点路径。

        Returns:
            若路径存在则返回 `True`，否则返回 `False`。

        Example:
            ```python
            if context.has_value("project.owner.name"):
                owner = context.get_value("project.owner.name")
            ```
        """

        return self.get_value(path, _MISSING) is not _MISSING

    @staticmethod
    def _resolve_part(value: Any, part: str, default: Any) -> Any:
        if value is None:
            return default
        if isinstance(value, dict):
            return value.get(part, default)
        if isinstance(value, (list, tuple)) and part.isdigit():
            index = int(part)
            if 0 <= index < len(value):
                return value[index]
            return default
        if hasattr(value, part):
            return getattr(value, part)
        if hasattr(value, "model_dump"):
            return value.model_dump().get(part, default)
        if hasattr(value, "dict") and callable(value.dict):
            return value.dict().get(part, default)
        return default


@dataclass
class EngineOptions:
    """控制渲染引擎行为的配置对象。

    Args:
        duplicate_key_policy: 重复占位块 key 的处理策略，支持 `error` 和 `broadcast`。
        missing_renderer_policy: 缺少 renderer 时的处理策略，支持 `error` 和 `warn`。
        enable_text_field_replace: 是否启用普通文本与表格单元格字段替换。
        text_field_pattern: 文本字段匹配正则，默认支持 `{{path.to.value}}`。
        text_field_replace_mode: 文本替换模式，当前版本仅支持 `plain`。
        strict: 是否在校验失败时尽早报错。
    """

    duplicate_key_policy: str = "broadcast"
    missing_renderer_policy: str = "error"
    enable_text_field_replace: bool = True
    text_field_pattern: str = r"\{\{([\w\.]+)\}\}"
    text_field_replace_mode: str = "plain"
    strict: bool = True


class Content:
    """所有渲染内容类型的基类。

    业务 renderer 不直接操作底层 PPT 对象，而是返回 ``Content`` 子类，由
    SDK 统一完成写回。
    """


@dataclass
class TextContent(Content):
    """文本占位块的渲染结果。

    Args:
        text: 要写回到目标文本 shape 的完整字符串。

    Example:
        ```python
        return TextContent(text="2026 Q1 经营分析")
        ```
    """

    text: str


@dataclass
class ImageContent(Content):
    """图片占位块的渲染结果。

    Args:
        image_path: 本地图片路径，SDK 会将其插入到占位块区域。

    Example:
        ```python
        return ImageContent(image_path="assets/logo.png")
        ```
    """

    image_path: str


@dataclass
class TableContent(Content):
    """整表替换用的表格渲染结果。

    Args:
        headers: 表头行；为空时表示只有数据行。元素可为字符串或 `cell(...)`。
        rows: 二维数组，每个子列表代表一行。元素可为字符串或 `cell(...)`。

    Example:
        ```python
        return TableContent(
            headers=["风险", "等级"],
            rows=[["现金流", "高"], ["履约", "中"]],
        )
        ```
    """

    headers: list[Any]
    rows: list[list[Any]]


@dataclass
class _CellValue:
    text: str
    color: Optional[str] = None
    font_name: Optional[str] = None
    font_size: Optional[Any] = 12
    bold: Optional[bool] = None
    italic: Optional[bool] = None
    underline: Optional[bool] = None
    append: bool = False


def cell(
    text: Any,
    *,
    color: Optional[str] = None,
    bold: Optional[bool] = None,
    italic: Optional[bool] = None,
    underline: Optional[bool] = None,
    font_name: Optional[str] = None,
    font_size: Optional[Any] = 12,
    append: bool = False,
) -> _CellValue:
    """创建一个带可选字体样式覆盖的表格 cell 值。

    这是表格 cell 的推荐写法。相比直接构造对象，`cell(...)` 更接近
    `set_cell_text()` 的使用心智：先给文本，再按需补颜色和常用字体样式。

    Args:
        text: 需要写入 cell 的文本。
        color: 6 位 RGB 十六进制字符串，例如 `FF0000`。
        bold: 是否加粗。
        italic: 是否斜体。
        underline: 是否下划线。
        font_name: 字体名称。
        font_size: 字号，单位与 `python-pptx` 的 `font.size` 一致。默认 `12`。
        append: 若为 `True`，则保留原 cell 文本并追加一个新的 run。

    Returns:
        可放入 `TableContent` 或 `TableCellsContent` 的 cell 值对象。

    Example:
        ```python
        return TableContent(
            headers=["风险", "等级"],
            rows=[["现金流", cell("高", color="FF0000", bold=True)]],
        )
        ```
    """

    return _CellValue(
        text=str(text),
        color=color,
        font_name=font_name,
        font_size=Pt(font_size) if isinstance(font_size, (int, float)) else font_size,
        bold=bold,
        italic=italic,
        underline=underline,
        append=append,
    )


@dataclass
class TableCellsContent(Content):
    """局部更新表格 cell 的渲染结果。

    该类型只适用于原生表格 placeholder，不适用于文本框型 table 区域。

    Args:
        cells: 需要更新的 cell 映射。key 为 `(row, col)` 的 `0-based`
            绝对坐标，value 可为字符串或 `cell(...)`。空字符串表示清空该 cell。

    Example:
        ```python
        return TableCellsContent(cells={(1, 0): "现金流", (1, 1): "高"})
        ```
    """

    cells: dict[tuple[int, int], Any]


@dataclass
class ChartContent(Content):
    """图表占位块的渲染结果。

    当前版本按图片方式写回，因此结构与 `ImageContent` 类似。

    Args:
        image_path: 图表导出的本地图片路径。
    """

    image_path: str


@dataclass
class Placeholder:
    """模板占位块的标准化描述对象。

    `Placeholder` 不是 `python-pptx` 原生类型，而是 SDK 在解析模板后生成的
    描述对象。它把 `shape.name = ph:<type>:<key>` 里的信息拆出来，供 renderer
    和写回逻辑统一使用。

    Args:
        type: 占位块类型，来自 `ph:<type>:<key>` 中的 `<type>`，例如 `text`、
            `image`、`table`、`chart`。
        key: 占位块业务 key，来自 `ph:<type>:<key>` 中的 `<key>`。
        slide_index: 占位块所在 slide 的 `0-based` 索引。
        shape_id: 原始 shape id，常用于日志与后续操作定位。
        shape_name: 原始 `shape.name`。
        left: 占位区域左侧坐标。
        top: 占位区域顶部坐标。
        width: 占位区域宽度。
        height: 占位区域高度。
        shape: 原始底层 shape 对象，仅供 SDK 内部写回使用。

    Example:
        如果模板里的 shape.name 是 `ph:text:title`，解析后会得到：

        - `placeholder.type == "text"`
        - `placeholder.key == "title"`
    """

    type: str
    key: str
    slide_index: int
    shape_id: int
    shape_name: str
    left: int
    top: int
    width: int
    height: int
    shape: Any


@dataclass
class RenderResult:
    """一次模板渲染的结构化返回结果。

    Args:
        success: 是否渲染成功。
        output_path: 若调用时传入了输出路径，则为最终落盘路径。
        output_bytes: 渲染后的 PPT 字节流。
        rendered_count: 实际完成 shape 级渲染的占位块数量。
        skipped_count: 预留字段，用于记录未处理数量。
        warnings: 渲染阶段产生的 warning，例如缺失文本字段。
    """

    success: bool
    output_path: Optional[str] = None
    output_bytes: Optional[bytes] = None
    rendered_count: int = 0
    skipped_count: int = 0
    warnings: list[str] = field(default_factory=list)


@dataclass
class ValidationReport:
    """静态模板校验结果。

    Args:
        success: 是否通过静态校验。
        placeholder_count: 模板中识别到的合法占位块数量。
        errors: 阻断型问题列表。
        warnings: 非阻断型问题列表。
        unused_renderers: 已注册但未使用的 renderer key。
    """

    success: bool
    placeholder_count: int = 0
    errors: list[str] = field(default_factory=list)
    warnings: list[str] = field(default_factory=list)
    unused_renderers: list[str] = field(default_factory=list)


@dataclass
class TextReplaceResult:
    """独立文本替换操作的返回结果。

    Args:
        replaced_count: 本次命中的字段总数。
        warnings: 替换过程中产生的 warning，例如缺失字段。
    """

    replaced_count: int = 0
    warnings: list[str] = field(default_factory=list)


class BaseRenderer:
    """自定义 renderer 的基类。

    业务侧通常继承该类并实现 `render()`，或使用 `RendererRegistry` 的函数
    式注册方式。一个 renderer 不直接操作 PPT shape，而是返回一个 `Content`
    子类，让 SDK 统一完成写回。

    `supported_types` 用于声明这个 renderer 能处理哪些 placeholder `type`。
    例如一个只负责文本的 renderer 可以声明 `{"text"}`。

    Example:
        ```python
        class TitleRenderer(BaseRenderer):
            supported_types = {"text"}

            def render(self, placeholder, context, **kwargs):
                return TextContent(text=context.get_value("project.name", "未命名项目"))
        ```
    """

    supported_types: set[str] = set()

    def render(self, placeholder: Placeholder, context: RenderContext, **kwargs):
        """根据占位块和上下文生成渲染结果。

        这是 renderer 的核心接口。调用时，SDK 已经根据 placeholder 的 `key`
        找到了对应 renderer，并把当前占位块和业务上下文传进来。你的任务只有一件事：
        返回与 `placeholder.type` 匹配的 `Content`。

        Args:
            placeholder: 当前占位块描述。
            context: 当前渲染上下文。
            **kwargs: 注册 renderer 时绑定的固定参数。

        Returns:
            `TextContent`、`ImageContent`、`TableContent`、`TableCellsContent`
            或 `ChartContent`。

        Example:
            ```python
            def render_risk_table(placeholder, context):
                if context.get_value("mode") == "patch":
                    return TableCellsContent(cells={(1, 0): cell("现金流", color="FF0000")})
                return TableContent(
                    headers=["风险", "等级"],
                    rows=[["现金流", "高"]],
                )
            ```
        """

        raise NotImplementedError


class _FunctionRenderer(BaseRenderer):
    def __init__(self, func: Callable, bound_kwargs: Optional[dict[str, Any]] = None):
        self._func = func
        self._bound_kwargs = bound_kwargs or {}

    def render(self, placeholder: Placeholder, context: RenderContext, **kwargs):
        merged_kwargs = {**self._bound_kwargs, **kwargs}
        return self._func(placeholder, context, **merged_kwargs)


class _BoundRenderer(BaseRenderer):
    def __init__(self, renderer: BaseRenderer, bound_kwargs: dict[str, Any]):
        self._renderer = renderer
        self._bound_kwargs = bound_kwargs
        self.supported_types = getattr(renderer, "supported_types", set())

    def render(self, placeholder: Placeholder, context: RenderContext, **kwargs):
        merged_kwargs = {**self._bound_kwargs, **kwargs}
        return self._renderer.render(placeholder, context, **merged_kwargs)


class RendererRegistry:
    """管理模板占位块 renderer 的注册与查询。

    它是模板 `key` 与业务渲染逻辑之间的桥梁。注意：注册是按 `key` 做的，
    不是按 `type` 做的。

    例如：
    - 模板里有 `ph:text:title`
    - 这里注册 `key="title"` 的 renderer
    - 运行时 SDK 会把这个 placeholder 分发给该 renderer

    `type` 的作用不是“选 renderer”，而是“校验 renderer 返回值是否正确”。
    例如 `ph:text:title` 对应的 renderer 必须返回 `TextContent`。

    Example:
        ```python
        registry = RendererRegistry()

        @registry.renderer("title")
        def render_title(placeholder, context):
            return TextContent(text="经营分析")
        ```
    """

    def __init__(self) -> None:
        self._renderers: dict[str, BaseRenderer] = {}

    def register(self, key: str, renderer: BaseRenderer, **bound_kwargs: Any) -> None:
        """注册类式 renderer。

        Args:
            key: 模板中的占位块 key，例如 `"title"`。
            renderer: `BaseRenderer` 实例。
            **bound_kwargs: 注册时绑定的固定参数，可用于同一个 renderer 适配不同 key。

        Example:
            ```python
            registry.register("title", TitleRenderer())
            registry.register("subtitle", TitleRenderer(), prefix="副标题：")
            ```
        """

        self._renderers[key] = _BoundRenderer(renderer, bound_kwargs) if bound_kwargs else renderer

    def register_func(self, key: str, func: Callable, **bound_kwargs: Any) -> None:
        """注册函数式 renderer。

        Args:
            key: 模板中的占位块 key。
            func: 接收 `(placeholder, context, **kwargs)` 并返回内容对象的函数。
            **bound_kwargs: 注册时绑定的固定参数，会在调用时传给函数。

        Example:
            ```python
            registry.register_func(
                "risk_table",
                lambda placeholder, context: TableCellsContent(cells={(1, 0): cell("现金流", color="FF0000")}),
            )
            ```
        """

        self.register(key, _FunctionRenderer(func, bound_kwargs))

    def renderer(self, key: str, **bound_kwargs: Any):
        """返回装饰器，用于以声明式方式注册 renderer。

        Args:
            key: 模板中的占位块 key。
            **bound_kwargs: 注册时绑定的固定参数。

        Returns:
            可直接装饰函数的注册器。

        Example:
            ```python
            @registry.renderer("title", prefix="主标题")
            def render_title(placeholder, context, prefix):
                return TextContent(text=f"{prefix}: 经营分析")
            ```
        """

        def decorator(func: Callable):
            self.register_func(key, func, **bound_kwargs)
            return func

        return decorator

    def get(self, key: str):
        """按 key 获取 renderer；若不存在则返回 ``None``。"""

        return self._renderers.get(key)

    def keys(self) -> list[str]:
        """返回已注册的全部 key，按字典序排序。"""

        return sorted(self._renderers.keys())


class PptxAdapter:
    """单文件版内部使用的 Presentation 读写适配层。

    该类不直接面向业务方，但它承载了模板加载、输出保存、shape 定位和内容写回
    的核心能力，因此单文件里保留了完整实现。
    """

    def load(self, template_path: Optional[str] = None, template_bytes: Optional[bytes] = None) -> Presentation:
        if bool(template_path) == bool(template_bytes):
            raise ValueError("exactly one of template_path or template_bytes must be provided")
        if template_path:
            return Presentation(template_path)
        from io import BytesIO

        return Presentation(BytesIO(template_bytes))

    @staticmethod
    def save_to_bytes(presentation) -> bytes:
        from io import BytesIO

        buffer = BytesIO()
        presentation.save(buffer)
        return buffer.getvalue()

    @staticmethod
    def save_to_path(presentation, output_path: str) -> None:
        presentation.save(output_path)

    @staticmethod
    def get_slide(presentation, slide_index: int):
        if slide_index < 0 or slide_index >= len(presentation.slides):
            raise ShapeOperationError(f"slide index {slide_index} out of range")
        return presentation.slides[slide_index]

    @staticmethod
    def find_shape(slide, shape_locator: Union[int, str]):
        if isinstance(shape_locator, int):
            for shape in slide.shapes:
                if getattr(shape, "shape_id", None) == shape_locator:
                    return shape
            raise ShapeOperationError(f"shape id {shape_locator} not found")
        for shape in slide.shapes:
            if getattr(shape, "name", None) == shape_locator:
                return shape
        raise ShapeOperationError(f"shape '{shape_locator}' not found")

    @staticmethod
    def write_content(presentation, placeholder: Placeholder, content: Content) -> None:
        slide = presentation.slides[placeholder.slide_index]
        shape = placeholder.shape
        if isinstance(content, TextContent):
            if not getattr(shape, "has_text_frame", False):
                raise ShapeOperationError(f"shape '{placeholder.shape_name}' cannot accept text content")
            PptxAdapter._set_text_frame_text_preserving_style(shape.text_frame, content.text)
            return
        if isinstance(content, (ImageContent, ChartContent)):
            slide.shapes.add_picture(content.image_path, placeholder.left, placeholder.top, placeholder.width, placeholder.height)
            PptxAdapter._remove_shape(shape)
            return
        if isinstance(content, TableContent):
            if getattr(shape, "has_table", False):
                PptxAdapter._rewrite_table(shape.table, content)
                return
            rows, cols, grid = PptxAdapter._build_table_grid(content)
            table_shape = slide.shapes.add_table(rows, cols, placeholder.left, placeholder.top, placeholder.width, placeholder.height)
            table = table_shape.table
            for row_index, row_values in enumerate(grid):
                for col_index, value in enumerate(row_values):
                    table.cell(row_index, col_index).text = value
            PptxAdapter._remove_shape(shape)
            return
        if isinstance(content, TableCellsContent):
            if not getattr(shape, "has_table", False):
                raise ShapeOperationError(
                    f"shape '{placeholder.shape_name}' cannot accept partial table cell updates"
                )
            PptxAdapter._patch_table_cells(shape.table, content.cells)
            return
        raise ShapeOperationError(f"unsupported content type '{type(content).__name__}'")

    @staticmethod
    def _build_table_grid(content: TableContent) -> tuple[int, int, list[list[Any]]]:
        grid: list[list[Any]] = []
        if content.headers:
            grid.append(list(content.headers))
        for row in content.rows:
            grid.append(list(row))
        cols = max((len(row) for row in grid), default=1)
        normalized = [row + [""] * (cols - len(row)) for row in grid] or [[""]]
        return len(normalized), cols, normalized

    @staticmethod
    def _rewrite_table(table, content: TableContent) -> bool:
        rows, cols, grid = PptxAdapter._build_table_grid(content)
        if len(table.rows) != rows or len(table.columns) != cols:
            raise ShapeOperationError(
                f"table placeholder size mismatch: template is {len(table.rows)}x{len(table.columns)}, "
                f"content is {rows}x{cols}"
            )
        for row_index, row_values in enumerate(grid):
            for col_index, value in enumerate(row_values):
                text, style_override = PptxAdapter._normalize_table_cell_value(value)
                PptxAdapter._write_table_cell_text(table.cell(row_index, col_index).text_frame, text, style_override)
        return True

    @staticmethod
    def _patch_table_cells(table, cells: dict[tuple[int, int], Any]) -> None:
        max_rows = len(table.rows)
        max_cols = len(table.columns)
        for coordinates, value in cells.items():
            if not isinstance(coordinates, tuple) or len(coordinates) != 2:
                raise ShapeOperationError("table cell coordinates must be (row, col) tuples")
            row_index, col_index = coordinates
            if not isinstance(row_index, int) or not isinstance(col_index, int):
                raise ShapeOperationError("table cell coordinates must use integer row and col indexes")
            if row_index < 0 or row_index >= max_rows or col_index < 0 or col_index >= max_cols:
                raise ShapeOperationError(
                    f"table cell coordinate ({row_index}, {col_index}) out of range for {max_rows}x{max_cols} table"
                )
            text, style_override = PptxAdapter._normalize_table_cell_value(value)
            PptxAdapter._write_table_cell_text(table.cell(row_index, col_index).text_frame, text, style_override)

    @staticmethod
    def _normalize_table_cell_value(value: Any) -> tuple[str, Optional[_CellValue]]:
        if isinstance(value, _CellValue):
            return str(value.text), value
        if isinstance(value, str):
            return value, None
        raise ShapeOperationError(
            f"table cell value must be str or cell(...), got '{type(value).__name__}'"
        )

    @staticmethod
    def _write_table_cell_text(text_frame, text: str, style_override: Optional[_CellValue] = None) -> None:
        if style_override is not None and style_override.append:
            PptxAdapter._append_run_to_text_frame(text_frame, text, style_override)
            return
        PptxAdapter._set_text_frame_text_preserving_style(text_frame, text, style_override)

    @staticmethod
    def _remove_shape(shape) -> None:
        parent = shape.element.getparent()
        if parent is None:
            raise ShapeOperationError("unable to remove placeholder shape after overlay insertion")
        parent.remove(shape.element)

    @staticmethod
    def _capture_text_style(text_frame) -> dict[str, Any]:
        paragraph = text_frame.paragraphs[0] if text_frame.paragraphs else None
        run = paragraph.runs[0] if paragraph and paragraph.runs else None
        style = {
            "alignment": getattr(paragraph, "alignment", None),
            "level": getattr(paragraph, "level", None),
            "line_spacing": getattr(paragraph, "line_spacing", None),
            "space_before": getattr(paragraph, "space_before", None),
            "space_after": getattr(paragraph, "space_after", None),
            "font_name": None,
            "font_size": None,
            "font_bold": None,
            "font_italic": None,
            "font_underline": None,
            "font_rgb": None,
        }
        if run is not None:
            font = run.font
            style["font_name"] = font.name
            style["font_size"] = font.size
            style["font_bold"] = font.bold
            style["font_italic"] = font.italic
            style["font_underline"] = font.underline
            try:
                style["font_rgb"] = font.color.rgb
            except Exception:
                style["font_rgb"] = None
        return style

    @staticmethod
    def _set_text_frame_text_preserving_style(
        text_frame,
        text: str,
        style_override: Optional[_CellValue] = None,
    ) -> None:
        style = PptxAdapter._capture_text_style(text_frame)
        text_frame.clear()
        paragraph = text_frame.paragraphs[0]
        if style["alignment"] is not None:
            paragraph.alignment = style["alignment"]
        if style["level"] is not None:
            paragraph.level = style["level"]
        if style["line_spacing"] is not None:
            paragraph.line_spacing = style["line_spacing"]
        if style["space_before"] is not None:
            paragraph.space_before = style["space_before"]
        if style["space_after"] is not None:
            paragraph.space_after = style["space_after"]
        run = paragraph.add_run()
        run.text = text
        font = run.font
        if style["font_name"] is not None:
            font.name = style["font_name"]
        if style["font_size"] is not None:
            font.size = style["font_size"]
        if style["font_bold"] is not None:
            font.bold = style["font_bold"]
        if style["font_italic"] is not None:
            font.italic = style["font_italic"]
        if style["font_underline"] is not None:
            font.underline = style["font_underline"]
        if style["font_rgb"] is not None:
            try:
                font.color.rgb = style["font_rgb"]
            except Exception:
                pass
        if style_override is not None:
            PptxAdapter._apply_text_style_override(font, style_override)

    @staticmethod
    def _append_run_to_text_frame(text_frame, text: str, style_override: _CellValue) -> None:
        paragraph = text_frame.paragraphs[0] if text_frame.paragraphs else text_frame.add_paragraph()
        run = paragraph.add_run()
        run.text = text
        PptxAdapter._apply_text_style_override(run.font, style_override)

    @staticmethod
    def _apply_text_style_override(font, style: _CellValue) -> None:
        if style.font_name is not None:
            font.name = style.font_name
        if style.font_size is not None:
            font.size = style.font_size
        if style.bold is not None:
            font.bold = style.bold
        if style.italic is not None:
            font.italic = style.italic
        if style.underline is not None:
            font.underline = style.underline
        if style.color is not None:
            if not re.fullmatch(r"[0-9A-Fa-f]{6}", style.color):
                raise ShapeOperationError("cell(..., color=...) must use a 6-digit RGB hex string")
            font.color.rgb = RGBColor.from_string(style.color.upper())


@dataclass
class ParsedTemplate:
    placeholders: list[Placeholder] = field(default_factory=list)
    invalid_placeholders: list[str] = field(default_factory=list)
    text_field_paths: set[str] = field(default_factory=set)


def iter_shapes(shape_collection) -> Iterable:
    """遍历 slide 中的 shape，并递归展开 group shape。"""

    for shape in shape_collection:
        yield shape
        if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
            yield from iter_shapes(shape.shapes)


def parse_presentation(presentation, text_field_pattern: str) -> ParsedTemplate:
    """扫描模板中的占位块和文本字段。"""

    parsed = ParsedTemplate()
    field_re = re.compile(text_field_pattern)
    for slide_index, slide in enumerate(presentation.slides):
        for shape in iter_shapes(slide.shapes):
            name = getattr(shape, "name", "") or ""
            if name.startswith("ph:"):
                match = PLACEHOLDER_RE.match(name)
                if not match:
                    parsed.invalid_placeholders.append(
                        f"slide {slide_index + 1} shape {getattr(shape, 'shape_id', '?')}: invalid placeholder name '{name}'"
                    )
                else:
                    parsed.placeholders.append(
                        Placeholder(
                            type=match.group(1),
                            key=match.group(2),
                            slide_index=slide_index,
                            shape_id=getattr(shape, "shape_id", -1),
                            shape_name=name,
                            left=getattr(shape, "left", 0),
                            top=getattr(shape, "top", 0),
                            width=getattr(shape, "width", 0),
                            height=getattr(shape, "height", 0),
                            shape=shape,
                        )
                    )
            if getattr(shape, "has_text_frame", False):
                parsed.text_field_paths.update(field_re.findall(shape.text_frame.text))
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    for cell in row.cells:
                        parsed.text_field_paths.update(field_re.findall(cell.text))
    return parsed


def validate_presentation(presentation, registry, options) -> ValidationReport:
    """对模板执行静态校验。"""

    parsed = parse_presentation(presentation, options.text_field_pattern)
    errors = list(parsed.invalid_placeholders)
    warnings: list[str] = []

    key_counter = Counter(placeholder.key for placeholder in parsed.placeholders)
    for key, count in sorted(key_counter.items()):
        if count < 2:
            continue
        message = f"duplicate placeholder key '{key}' appears {count} times"
        if options.duplicate_key_policy == "error":
            errors.append(message)
        else:
            warnings.append(message)

    used_keys = set()
    for placeholder in parsed.placeholders:
        renderer = registry.get(placeholder.key)
        if renderer is None:
            message = f"missing renderer for placeholder '{placeholder.shape_name}' on slide {placeholder.slide_index + 1}"
            if options.missing_renderer_policy == "warn":
                warnings.append(message)
            else:
                errors.append(message)
            continue
        used_keys.add(placeholder.key)
        supported_types = getattr(renderer, "supported_types", set()) or set()
        if supported_types and placeholder.type not in supported_types:
            errors.append(f"renderer '{placeholder.key}' does not support placeholder type '{placeholder.type}'")

    unused_renderers = sorted(set(registry.keys()) - used_keys)
    if unused_renderers:
        warnings.append(f"unused renderers: {', '.join(unused_renderers)}")

    return ValidationReport(
        success=not errors,
        placeholder_count=len(parsed.placeholders),
        errors=errors,
        warnings=warnings,
        unused_renderers=unused_renderers,
    )


class TextReplacer:
    """执行普通文本框和表格单元格中的字段替换。

    适合在不走 shape 级 renderer 的场景下，单独对现有 PPT 做字段替换。

    Args:
        pattern: 可选字段匹配正则。未传时使用默认 `{{field.path}}` 形式。

    Example:
        ```python
        from pptx import Presentation

        prs = Presentation("text_replace_template.pptx")
        result = TextReplacer().replace_presentation_text(
            prs,
            context=RenderContext(data={"project": {"name": "Aurora"}}),
        )
        ```
    """

    def __init__(self, pattern: Optional[str] = None):
        self.pattern = pattern

    def replace_presentation_text(
        self,
        presentation,
        context: RenderContext,
        rendered_shape_ids: Optional[set[int]] = None,
        pattern: Optional[str] = None,
    ) -> TextReplaceResult:
        """对整个 Presentation 执行字段替换。

        Args:
            presentation: `python-pptx` 的 `Presentation` 实例。
            context: 字段取值上下文。
            rendered_shape_ids: 可选 shape id 集合；这些 shape 会被跳过。
            pattern: 本次调用临时覆盖的匹配正则。

        Returns:
            `TextReplaceResult`，包含替换次数与 warning 列表。

        Raises:
            FieldReplaceError: 替换过程发生异常时抛出。

        Example:
            ```python
            prs = Presentation("template.pptx")
            result = TextReplacer().replace_presentation_text(
                prs,
                context=RenderContext(data={"project": {"name": "Aurora"}}),
            )
            ```
        """

        field_re = re.compile(pattern or self.pattern or r"\{\{([\w\.]+)\}\}")
        rendered_shape_ids = rendered_shape_ids or set()
        warnings: list[str] = []
        replaced_count = 0
        try:
            for slide in presentation.slides:
                for shape in iter_shapes(slide.shapes):
                    if getattr(shape, "shape_id", None) in rendered_shape_ids:
                        continue
                    if getattr(shape, "has_text_frame", False):
                        new_text, count = self._replace_text(shape.text_frame.text, field_re, context, warnings)
                        shape.text = new_text
                        replaced_count += count
                    if getattr(shape, "has_table", False):
                        for row in shape.table.rows:
                            for cell in row.cells:
                                new_text, count = self._replace_text(cell.text, field_re, context, warnings)
                                cell.text = new_text
                                replaced_count += count
        except Exception as exc:
            raise FieldReplaceError(str(exc)) from exc

        return TextReplaceResult(replaced_count=replaced_count, warnings=warnings)

    @staticmethod
    def _replace_text(text: str, field_re, context: RenderContext, warnings: list[str]) -> tuple[str, int]:
        replacements = 0

        def repl(match):
            nonlocal replacements
            path = match.group(1)
            value = context.get_value(path)
            replacements += 1
            if value is None:
                warnings.append(f"missing text field '{path}'")
                return ""
            return str(value)

        return field_re.sub(repl, text), replacements


CONTENT_TYPE_MAP = {
    "text": (TextContent,),
    "image": (ImageContent,),
    "table": (TableContent, TableCellsContent),
    "chart": (ChartContent,),
}


class PptTemplateEngine:
    """PPT 模板渲染主引擎。

    这是单文件版最主要的入口。它负责把“模板里的 placeholder”与“代码里的 renderer”
    连接起来，并在最后输出渲染后的 PPT。

    Args:
        registry: 已注册 renderer 的 `RendererRegistry`。
        options: 可选引擎配置；未传时使用默认 `EngineOptions()`。

    Example:
        ```python
        engine = PptTemplateEngine(registry=registry)
        result = engine.render(
            template_path="report_template.pptx",
            output_path="report_output.pptx",
            context=context,
        )
        ```
    """

    def __init__(self, registry, options: Optional[EngineOptions] = None):
        self.registry = registry
        self.options = options or EngineOptions()
        self.adapter = PptxAdapter()
        self.text_replacer = TextReplacer(self.options.text_field_pattern)

    def render(
        self,
        template_path: Optional[str] = None,
        template_bytes: Optional[bytes] = None,
        output_path: Optional[str] = None,
        context: Optional[RenderContext] = None,
        operations_builder: Optional[Callable[[Any, RenderContext], None]] = None,
    ) -> RenderResult:
        """执行模板渲染。

        整个流程如下：
        1. 读取模板并扫描所有 `ph:<type>:<key>` placeholder。
        2. 按 `key` 去 `RendererRegistry` 查找 renderer。
        3. 调用 renderer，拿到 `Content`。
        4. 校验 `Content` 是否和 placeholder 的 `type` 匹配。
        5. 把内容写回 PPT。
        6. 最后再做普通文本和表格单元格中的 `{{field}}` 替换。

        如果你第一次接触这个库，可以把它理解成：
        “模板负责标记位置，renderer 负责给数据，engine 负责把两者拼起来”。

        Args:
            template_path: 模板文件路径，与 `template_bytes` 二选一。
            template_bytes: 模板字节流，与 `template_path` 二选一。
            output_path: 可选输出路径；传入后会在返回 bytes 的同时落盘。
            context: 渲染上下文；为空时使用空上下文。
            operations_builder: 可选后处理回调。回调会在渲染和文本替换完成后
                收到 `PptOperations` 与 `RenderContext`，用于根据外部参数决定删页、
                插页等结构操作。

        Returns:
            `RenderResult`，包含输出 bytes、渲染计数和 warning。

        Raises:
            PlaceholderFormatError: 模板中存在非法占位块命名。
            DuplicatePlaceholderError: 重复 key 且策略为 `error`。
            RendererNotFoundError: 模板存在未注册 renderer 的占位块，且策略为 `error`。
            ContentTypeMismatchError: renderer 返回类型与占位块类型不匹配。

        Example:
            ```python
            result = engine.render(
                template_path="template.pptx",
                output_path="output.pptx",
                context=RenderContext(data={"project": {"name": "Aurora"}}),
            )
            ```
        """

        context = context or RenderContext(data={})
        presentation = self.adapter.load(template_path=template_path, template_bytes=template_bytes)
        parsed = parse_presentation(presentation, self.options.text_field_pattern)
        if parsed.invalid_placeholders:
            raise PlaceholderFormatError(parsed.invalid_placeholders[0])

        placeholder_groups = defaultdict(list)
        for placeholder in parsed.placeholders:
            placeholder_groups[placeholder.key].append(placeholder)

        rendered_shape_ids: set[int] = set()
        rendered_count = 0
        skipped_count = 0
        warnings: list[str] = []

        for key, placeholders in placeholder_groups.items():
            if len(placeholders) > 1 and self.options.duplicate_key_policy == "error":
                raise DuplicatePlaceholderError(f"duplicate placeholder key '{key}' appears {len(placeholders)} times")
            renderer = self.registry.get(key)
            if renderer is None:
                if self.options.missing_renderer_policy == "warn":
                    warnings.append(f"missing renderer for placeholder key '{key}'")
                    skipped_count += len(placeholders)
                    continue
                raise RendererNotFoundError(f"missing renderer for placeholder key '{key}'")
            content = renderer.render(placeholders[0], context)
            expected_types = CONTENT_TYPE_MAP[placeholders[0].type]
            if not isinstance(content, expected_types):
                expected_names = ", ".join(expected_type.__name__ for expected_type in expected_types)
                raise ContentTypeMismatchError(
                    f"renderer '{key}' returned {type(content).__name__}, expected one of: {expected_names}"
                )
            for placeholder in placeholders:
                self.adapter.write_content(presentation, placeholder, content)
                rendered_shape_ids.add(placeholder.shape_id)
                rendered_count += 1

        if self.options.enable_text_field_replace:
            replace_result = self.text_replacer.replace_presentation_text(
                presentation,
                context=context,
                rendered_shape_ids=rendered_shape_ids,
                pattern=self.options.text_field_pattern,
            )
            warnings.extend(replace_result.warnings)

        if operations_builder is not None:
            operations = PptOperations(presentation, adapter=self.adapter)
            operations_builder(operations, context)

        output_bytes = self.adapter.save_to_bytes(presentation)
        if output_path:
            self.adapter.save_to_path(presentation, output_path)

        return RenderResult(
            success=True,
            output_path=output_path,
            output_bytes=output_bytes,
            rendered_count=rendered_count,
            skipped_count=skipped_count,
            warnings=warnings,
        )

    def validate(
        self,
        template_path: Optional[str] = None,
        template_bytes: Optional[bytes] = None,
    ) -> ValidationReport:
        """对模板执行静态校验。

        Args:
            template_path: 模板文件路径，与 `template_bytes` 二选一。
            template_bytes: 模板字节流，与 `template_path` 二选一。

        Returns:
            `ValidationReport`，用于查看错误、warning 和未使用 renderer。

        Example:
            ```python
            report = engine.validate(template_path="template.pptx")
            ```
        """

        presentation = self.adapter.load(template_path=template_path, template_bytes=template_bytes)
        return validate_presentation(presentation, self.registry, self.options)


def _tag(namespace: str, name: str) -> str:
    return f"{{{namespace}}}{name}"


class PptOperations:
    """封装常见的 PPT 结构与表格操作。

    适合在渲染前后对文档做结构调整，例如插页、删页、按 section 组织页面，
    或对表格执行删行删列与合并操作。所有公开索引都是 ``0-based``。

    Args:
        presentation: 已加载的 `Presentation` 实例。
        adapter: 可选适配层；默认使用 `PptxAdapter`。

    Example:
        ```python
        ops = PptOperations.load(template_path="operations_template.pptx")
        ops.insert_slide(target_index=1, layout_index=6)
        ops.add_section(name="正文", start_slide_index=1)
        ops.save_to_path("operations_output.pptx")
        ```
    """

    def __init__(self, presentation, adapter: Optional[PptxAdapter] = None):
        self.presentation = presentation
        self.adapter = adapter or PptxAdapter()

    @classmethod
    def load(cls, template_path: Optional[str] = None, template_bytes: Optional[bytes] = None):
        """从路径或字节流加载 PPT 并创建操作对象。

        Args:
            template_path: 模板路径，与 `template_bytes` 二选一。
            template_bytes: 模板字节流，与 `template_path` 二选一。

        Returns:
            `PptOperations` 实例。

        Example:
            ```python
            ops = PptOperations.load(template_path="report.pptx")
            ```
        """

        adapter = PptxAdapter()
        return cls(adapter.load(template_path=template_path, template_bytes=template_bytes), adapter=adapter)

    def save_to_bytes(self) -> bytes:
        """将当前 Presentation 保存为内存字节流。

        Returns:
            当前 PPT 的二进制内容。

        Example:
            ```python
            payload = ops.save_to_bytes()
            ```
        """

        return self.adapter.save_to_bytes(self.presentation)

    def save_to_path(self, output_path: str) -> None:
        """将当前 Presentation 保存到指定路径。

        Args:
            output_path: 输出文件路径。

        Example:
            ```python
            ops.save_to_path("output.pptx")
            ```
        """

        self.adapter.save_to_path(self.presentation, output_path)

    def delete_slide(self, slide_index: int) -> int:
        """删除指定索引的 slide。

        Args:
            slide_index: `0-based` slide 索引。

        Returns:
            被删除 slide 的内部 `slide_id`。
        """

        slide_id_el = self._get_slide_id_element(slide_index)
        slide_id = int(slide_id_el.get("id"))
        rel_id = slide_id_el.get(_tag(R_NS, "id"))
        self.presentation.part.drop_rel(rel_id)
        self.presentation.slides._sldIdLst.remove(slide_id_el)
        if self._has_sections():
            groups = self._read_sections()
            for group in groups:
                group["slide_ids"] = [item for item in group["slide_ids"] if item != slide_id]
            self._write_sections([group for group in groups if group["slide_ids"]])
        return slide_id

    def insert_slide(self, target_index: int, layout_index: int):
        """使用模板现有 `layout_index` 新建并插入 slide。

        Args:
            target_index: 新 slide 的插入位置，使用 `0-based` 索引。
            layout_index: `presentation.slide_layouts` 中的 layout 索引。

        Returns:
            新建的 slide 对象。
        """

        if target_index < 0 or target_index > len(self.presentation.slides):
            raise OperationError(f"target_index {target_index} out of range")
        if layout_index < 0 or layout_index >= len(self.presentation.slide_layouts):
            raise OperationError(f"layout_index {layout_index} out of range")

        previous_slide_ids = self._slide_ids_in_order()
        groups = self._read_sections()
        slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[layout_index])
        slide_id = slide.slide_id
        slide_id_el = self.presentation.slides._sldIdLst[-1]
        self.presentation.slides._sldIdLst.remove(slide_id_el)
        self.presentation.slides._sldIdLst.insert(target_index, slide_id_el)

        if groups:
            if target_index >= len(previous_slide_ids):
                target_group = len(groups) - 1
                groups[target_group]["slide_ids"].append(slide_id)
            else:
                anchor_slide_id = previous_slide_ids[target_index]
                target_group = self._group_index_for_slide(groups, anchor_slide_id)
                anchor_pos = groups[target_group]["slide_ids"].index(anchor_slide_id)
                groups[target_group]["slide_ids"].insert(anchor_pos, slide_id)
            self._write_sections(groups)
        return slide

    def add_section(self, name: str, start_slide_index: int) -> None:
        """在指定 slide 位置开始一个新的 section。

        Args:
            name: section 名称。
            start_slide_index: section 起始 slide 的 `0-based` 索引。

        Notes:
            若模板当前没有 section，会自动初始化 section 列表。
            若目标 slide 已经是某个 section 的起始页，则该 section 会被重命名。
        """

        slide_ids = self._slide_ids_in_order()
        if not slide_ids:
            raise OperationError("cannot add a section to an empty presentation")
        if start_slide_index < 0 or start_slide_index >= len(slide_ids):
            raise OperationError(f"start_slide_index {start_slide_index} out of range")

        start_slide_id = slide_ids[start_slide_index]
        groups = self._read_sections()
        if not groups:
            before = slide_ids[:start_slide_index]
            after = slide_ids[start_slide_index:]
            groups = []
            if before:
                groups.append(self._make_section(DEFAULT_SECTION_NAME, before))
            groups.append(self._make_section(name, after))
            self._write_sections(groups)
            return

        group_index = self._group_index_for_slide(groups, start_slide_id)
        group = groups[group_index]
        first_slide_id = group["slide_ids"][0]
        if first_slide_id == start_slide_id:
            group["name"] = name
            self._write_sections(groups)
            return

        split_at = group["slide_ids"].index(start_slide_id)
        before = group["slide_ids"][:split_at]
        after = group["slide_ids"][split_at:]
        groups[group_index] = self._make_section(group["name"], before, guid=group["id"])
        groups.insert(group_index + 1, self._make_section(name, after))
        self._write_sections(groups)

    def delete_section(self, section_index: int) -> None:
        """删除指定 section，但保留其中 slides。

        Args:
            section_index: `0-based` section 索引。

        Notes:
            被删除 section 中的 slides 会并入相邻 section，不会从文档中删除。
        """

        groups = self._read_sections()
        if not groups:
            raise OperationError("presentation does not contain sections")
        if section_index < 0 or section_index >= len(groups):
            raise OperationError(f"section_index {section_index} out of range")
        if len(groups) == 1:
            self._write_sections([])
            return

        removed = groups.pop(section_index)
        if section_index == 0:
            groups[0]["slide_ids"] = removed["slide_ids"] + groups[0]["slide_ids"]
        else:
            groups[section_index - 1]["slide_ids"].extend(removed["slide_ids"])
        self._write_sections(groups)

    def delete_table_row(self, slide_index: int, shape_locator: Union[int, str], row_index: int) -> None:
        """删除指定表格中的一行。

        Args:
            slide_index: 目标 slide 的 `0-based` 索引。
            shape_locator: 推荐传 `shape_id`，也支持 `shape_name`。
            row_index: 待删除的 `0-based` 行索引。

        Raises:
            OperationError: 行越界，或表格已存在合并单元格。
            ShapeOperationError: 目标 shape 不存在或不是表格。
        """

        table = self._resolve_table(slide_index, shape_locator)
        self._ensure_unmerged_table(table)
        if row_index < 0 or row_index >= len(table.rows):
            raise OperationError(f"row_index {row_index} out of range")
        table._tbl.remove(table._tbl.tr_lst[row_index])

    def delete_table_column(self, slide_index: int, shape_locator: Union[int, str], column_index: int) -> None:
        """删除指定表格中的一列。

        Args:
            slide_index: 目标 slide 的 `0-based` 索引。
            shape_locator: 推荐传 `shape_id`，也支持 `shape_name`。
            column_index: 待删除的 `0-based` 列索引。
        """

        table = self._resolve_table(slide_index, shape_locator)
        self._ensure_unmerged_table(table)
        if column_index < 0 or column_index >= len(table.columns):
            raise OperationError(f"column_index {column_index} out of range")
        table._tbl.tblGrid.remove(table._tbl.tblGrid.gridCol_lst[column_index])
        for tr in table._tbl.tr_lst:
            tr.remove(tr.tc_lst[column_index])

    def patch_table_cells(
        self,
        slide_index: int,
        shape_locator: Union[int, str],
        cells: dict[tuple[int, int], str],
    ) -> None:
        """只更新指定表格单元格的文本。

        Args:
            slide_index: 目标 slide 的 `0-based` 索引。
            shape_locator: 推荐传 `shape_id`，也支持 `shape_name`。
            cells: 需要更新的 cell 文本映射，key 为 `(row, col)` 坐标。

        Raises:
            ShapeOperationError: 目标 shape 不是表格，或坐标越界时抛出。

        Example:
            ```python
            ops.patch_table_cells(
                slide_index=0,
                shape_locator="risk-table",
                cells={(1, 0): "现金流", (1, 1): "高"},
            )
            ```
        """

        table = self._resolve_table(slide_index, shape_locator)
        self.adapter._patch_table_cells(table, cells)

    def merge_table_cells(
        self,
        slide_index: int,
        shape_locator: Union[int, str],
        first_row: int,
        first_col: int,
        last_row: int,
        last_col: int,
    ) -> None:
        """合并指定矩形区域内的表格单元格。

        Args:
            slide_index: 目标 slide 的 `0-based` 索引。
            shape_locator: 推荐传 `shape_id`，也支持 `shape_name`。
            first_row: 合并区域左上角行索引。
            first_col: 合并区域左上角列索引。
            last_row: 合并区域右下角行索引。
            last_col: 合并区域右下角列索引。

        Raises:
            OperationError: 合并区域非法或越界。
        """

        table = self._resolve_table(slide_index, shape_locator)
        self._validate_merge_bounds(table, first_row, first_col, last_row, last_col)
        table.cell(first_row, first_col).merge(table.cell(last_row, last_col))

    def _resolve_table(self, slide_index: int, shape_locator: Union[int, str]):
        slide = self.adapter.get_slide(self.presentation, slide_index)
        shape = self.adapter.find_shape(slide, shape_locator)
        if not getattr(shape, "has_table", False):
            raise ShapeOperationError("target shape is not a table")
        return shape.table

    @staticmethod
    def _validate_merge_bounds(table, first_row: int, first_col: int, last_row: int, last_col: int) -> None:
        if first_row > last_row or first_col > last_col:
            raise OperationError("merge bounds must define a top-left to bottom-right rectangle")
        if first_row < 0 or first_col < 0 or last_row >= len(table.rows) or last_col >= len(table.columns):
            raise OperationError("merge bounds out of range")

    @staticmethod
    def _ensure_unmerged_table(table) -> None:
        for tr in table._tbl.tr_lst:
            for tc in tr.tc_lst:
                if tc.get("rowSpan") or tc.get("gridSpan") or tc.get("hMerge") or tc.get("vMerge"):
                    raise OperationError("row/column deletion is not supported on merged tables")

    def _get_slide_id_element(self, slide_index: int):
        slide_ids = list(self.presentation.slides._sldIdLst)
        if slide_index < 0 or slide_index >= len(slide_ids):
            raise OperationError(f"slide_index {slide_index} out of range")
        return slide_ids[slide_index]

    def _slide_ids_in_order(self) -> list[int]:
        return [slide.slide_id for slide in self.presentation.slides]

    def _group_index_for_slide(self, groups: list[dict], slide_id: int) -> int:
        for index, group in enumerate(groups):
            if slide_id in group["slide_ids"]:
                return index
        raise OperationError(f"slide id {slide_id} is not assigned to any section")

    @staticmethod
    def _make_section(name: str, slide_ids: list[int], guid: Optional[str] = None) -> dict:
        return {"name": name, "id": guid or f"{{{str(uuid.uuid4()).upper()}}}", "slide_ids": slide_ids}

    def _has_sections(self) -> bool:
        return self._find_section_ext() is not None

    def _read_sections(self) -> list[dict]:
        section_ext = self._find_section_ext()
        if section_ext is None:
            return []
        section_lst = section_ext.find(_tag(P14_NS, "sectionLst"))
        if section_lst is None:
            return []
        groups = []
        for section in section_lst.findall(_tag(P14_NS, "section")):
            slide_ids = [int(sld_id.get("id")) for sld_id in section.find(_tag(P14_NS, "sldIdLst")).findall(_tag(P14_NS, "sldId"))]
            groups.append(
                {
                    "name": section.get("name") or DEFAULT_SECTION_NAME,
                    "id": section.get("id") or f"{{{str(uuid.uuid4()).upper()}}}",
                    "slide_ids": slide_ids,
                }
            )
        return groups

    def _write_sections(self, groups: list[dict]) -> None:
        presentation_el = self.presentation.part._element
        ext_lst = presentation_el.find(_tag(P_NS, "extLst"))
        section_ext = self._find_section_ext()
        if not groups:
            if section_ext is not None:
                ext_lst.remove(section_ext)
            if ext_lst is not None and len(ext_lst) == 0:
                presentation_el.remove(ext_lst)
            return

        ordered_ids = self._slide_ids_in_order()
        order_map = {slide_id: index for index, slide_id in enumerate(ordered_ids)}
        normalized_groups = []
        for group in groups:
            slide_ids = sorted([slide_id for slide_id in group["slide_ids"] if slide_id in order_map], key=order_map.__getitem__)
            if slide_ids:
                normalized_groups.append({**group, "slide_ids": slide_ids})

        if ext_lst is None:
            ext_lst = etree.SubElement(presentation_el, _tag(P_NS, "extLst"))
        if section_ext is None:
            section_ext = etree.SubElement(ext_lst, _tag(P_NS, "ext"), uri=SECTION_EXT_URI)
        else:
            for child in list(section_ext):
                section_ext.remove(child)

        section_lst = etree.SubElement(section_ext, _tag(P14_NS, "sectionLst"), nsmap={"p14": P14_NS})
        for group in normalized_groups:
            section_el = etree.SubElement(section_lst, _tag(P14_NS, "section"), name=group["name"], id=group["id"])
            slide_id_lst = etree.SubElement(section_el, _tag(P14_NS, "sldIdLst"))
            for slide_id in group["slide_ids"]:
                etree.SubElement(slide_id_lst, _tag(P14_NS, "sldId"), id=str(slide_id))

    def _find_section_ext(self):
        presentation_el = self.presentation.part._element
        ext_lst = presentation_el.find(_tag(P_NS, "extLst"))
        if ext_lst is None:
            return None
        for ext in ext_lst.findall(_tag(P_NS, "ext")):
            if ext.get("uri") == SECTION_EXT_URI:
                return ext
        return None


__all__ = [
    "BaseRenderer",
    "ChartContent",
    "Content",
    "ContentTypeMismatchError",
    "cell",
    "DuplicatePlaceholderError",
    "EngineOptions",
    "FieldReplaceError",
    "ImageContent",
    "OperationError",
    "Placeholder",
    "PlaceholderFormatError",
    "PptTemplateEngine",
    "PptOperations",
    "PptTemplateSdkError",
    "RenderContext",
    "RenderResult",
    "RendererNotFoundError",
    "RendererRegistry",
    "ShapeOperationError",
    "TableContent",
    "TableCellsContent",
    "TemplateParseError",
    "TextReplaceResult",
    "TextReplacer",
    "TextContent",
    "ValidationReport",
]

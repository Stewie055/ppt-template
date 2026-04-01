from __future__ import annotations

import base64
from dataclasses import dataclass
import importlib.util
from io import BytesIO
from pathlib import Path
import sys
from zipfile import ZipFile

import pytest
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches
from pptx.util import Pt

PROJECT_ROOT = Path(__file__).resolve().parents[1]
if str(PROJECT_ROOT / "singlefile") not in sys.path:
    sys.path.insert(0, str(PROJECT_ROOT / "singlefile"))

from ppt_template_sdk import (
    BaseRenderer,
    ChartContent,
    ContentTypeMismatchError,
    DuplicatePlaceholderError,
    cell,
    EngineOptions,
    ImageContent,
    OperationError,
    PptTemplateEngine,
    PptOperations,
    RenderContext,
    RendererRegistry,
    TableCellsContent,
    TableContent,
    TextReplacer,
    TextContent,
)


PNG_1X1 = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAusB9Y9n6n8AAAAASUVORK5CYII="
)


def _load_singlefile_module():
    module_path = PROJECT_ROOT / "singlefile" / "ppt_template_sdk.py"
    spec = importlib.util.spec_from_file_location("ppt_template_sdk_single", module_path)
    module = importlib.util.module_from_spec(spec)
    assert spec.loader is not None
    sys.modules[spec.name] = module
    spec.loader.exec_module(module)
    return module


def _write_png(path: Path) -> None:
    path.write_bytes(PNG_1X1)


def _build_template(path: Path) -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    title.name = "ph:text:title"
    title.text = "placeholder"

    logo = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(2), Inches(1))
    logo.name = "ph:image:logo"
    logo.text = "img"

    chart = slide.shapes.add_textbox(Inches(4), Inches(2), Inches(2), Inches(1))
    chart.name = "ph:chart:sales"
    chart.text = "chart"

    table = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(5), Inches(1.5))
    table.name = "ph:table:risk_table"
    table.text = "table"

    field_box = slide.shapes.add_textbox(Inches(1), Inches(5.3), Inches(4), Inches(1))
    field_box.text = "项目：{{project.name}} 日期：{{report_date}} 缺失：{{missing.value}}"

    native_table = slide.shapes.add_table(2, 2, Inches(1), Inches(6.1), Inches(5), Inches(1.5)).table
    native_table.cell(0, 0).text = "负责人"
    native_table.cell(0, 1).text = "{{owner.name}}"
    native_table.cell(1, 0).text = "项目"
    native_table.cell(1, 1).text = "{{project.name}}"

    prs.save(path)


def _build_dual_text_template(path: Path) -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    first = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    first.name = "ph:text:title"
    first.text = "first"
    second = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(4), Inches(1))
    second.name = "ph:text:subtitle"
    second.text = "second"
    prs.save(path)


def _build_styled_text_template(path: Path) -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    shape.name = "ph:text:title"
    paragraph = shape.text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = "placeholder"
    run.font.name = "Arial"
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.italic = True
    run.font.color.rgb = RGBColor(0x11, 0x22, 0x33)
    prs.save(path)


def _build_native_table_placeholder_template(path: Path) -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = slide.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(5), Inches(2))
    shape.name = "ph:table:risk_table"
    table = shape.table
    table.columns[0].width = Inches(3)
    table.columns[1].width = Inches(2)
    table.rows[0].height = Inches(0.8)
    table.rows[1].height = Inches(1.2)
    paragraph = table.cell(0, 0).text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = "风险"
    run.font.name = "Arial"
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x44, 0x55, 0x66)
    table.cell(0, 1).text = "等级"
    table.cell(1, 0).text = "旧值"
    table.cell(1, 1).text = "旧值"
    prs.save(path)


@dataclass
class Report:
    title: str


class TitleRenderer(BaseRenderer):
    supported_types = {"text"}

    def render(self, placeholder, context):
        report = context.extras["report"]
        return TextContent(text=report.title)


def test_render_full_flow(tmp_path: Path):
    template_path = tmp_path / "template.pptx"
    output_path = tmp_path / "output.pptx"
    logo_path = tmp_path / "logo.png"
    chart_path = tmp_path / "chart.png"
    _write_png(logo_path)
    _write_png(chart_path)
    _build_template(template_path)

    registry = RendererRegistry()
    registry.register("title", TitleRenderer())
    registry.register_func("logo", lambda placeholder, context: ImageContent(image_path=str(logo_path)))
    registry.register_func("sales", lambda placeholder, context: ChartContent(image_path=str(chart_path)))
    registry.register_func(
        "risk_table",
        lambda placeholder, context: TableContent(headers=["风险", "等级"], rows=[["现金流", "高"], ["履约", "中"]]),
    )

    engine = PptTemplateEngine(registry, EngineOptions(duplicate_key_policy="broadcast"))
    result = engine.render(
        template_path=str(template_path),
        output_path=str(output_path),
        context=RenderContext(
            data={"project": {"name": "北极星"}, "report_date": "2026-03-31", "owner": {"name": "Alice"}},
            extras={"report": Report(title="Q1 经营分析")},
        ),
    )

    assert result.success is True
    assert result.output_bytes
    assert output_path.exists()
    assert result.rendered_count == 4
    assert any("missing text field 'missing.value'" in warning for warning in result.warnings)

    rendered = Presentation(str(output_path))
    texts = [shape.text for shape in rendered.slides[0].shapes if getattr(shape, "has_text_frame", False)]
    assert "Q1 经营分析" in texts
    assert "项目：北极星 日期：2026-03-31 缺失：" in texts
    assert any(getattr(shape, "has_table", False) for shape in rendered.slides[0].shapes)


def test_render_from_bytes(tmp_path: Path):
    template_path = tmp_path / "template.pptx"
    _build_template(template_path)
    logo_path = tmp_path / "logo.png"
    chart_path = tmp_path / "chart.png"
    _write_png(logo_path)
    _write_png(chart_path)

    template_bytes = template_path.read_bytes()
    registry = RendererRegistry()
    registry.register_func("title", lambda placeholder, context: TextContent(text="Bytes"))
    registry.register_func("logo", lambda placeholder, context: ImageContent(image_path=str(logo_path)))
    registry.register_func("sales", lambda placeholder, context: ChartContent(image_path=str(chart_path)))
    registry.register_func("risk_table", lambda placeholder, context: TableContent(headers=["A"], rows=[["B"]]))

    engine = PptTemplateEngine(registry)
    result = engine.render(template_bytes=template_bytes, context=RenderContext(data={}))

    assert result.success is True
    assert isinstance(result.output_bytes, bytes)


def test_validate_reports_static_issues(tmp_path: Path):
    template_path = tmp_path / "template.pptx"
    _build_template(template_path)
    registry = RendererRegistry()
    registry.register_func("unused", lambda placeholder, context: TextContent(text="x"))

    engine = PptTemplateEngine(registry)
    report = engine.validate(template_path=str(template_path))

    assert report.success is False
    assert report.placeholder_count == 4
    assert any("missing renderer" in error for error in report.errors)
    assert report.unused_renderers == ["unused"]


def test_validate_missing_renderer_can_be_warning(tmp_path: Path):
    template_path = tmp_path / "template.pptx"
    _build_template(template_path)
    registry = RendererRegistry()
    registry.register_func("unused", lambda placeholder, context: TextContent(text="x"))

    engine = PptTemplateEngine(registry, EngineOptions(missing_renderer_policy="warn"))
    report = engine.validate(template_path=str(template_path))

    assert report.success is True
    assert report.errors == []
    assert any("missing renderer" in warning for warning in report.warnings)
    assert report.unused_renderers == ["unused"]


def test_duplicate_key_error_policy(tmp_path: Path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    first = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
    first.name = "ph:text:title"
    second = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(3), Inches(1))
    second.name = "ph:text:title"
    path = tmp_path / "dup.pptx"
    prs.save(path)

    registry = RendererRegistry()
    registry.register_func("title", lambda placeholder, context: TextContent(text="dup"))
    engine = PptTemplateEngine(registry, EngineOptions(duplicate_key_policy="error"))

    with pytest.raises(DuplicatePlaceholderError):
        engine.render(template_path=str(path), context=RenderContext(data={}))


def test_content_type_mismatch(tmp_path: Path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
    shape.name = "ph:image:cover"
    path = tmp_path / "mismatch.pptx"
    prs.save(path)

    registry = RendererRegistry()
    registry.register_func("cover", lambda placeholder, context: TextContent(text="wrong"))
    engine = PptTemplateEngine(registry)

    with pytest.raises(ContentTypeMismatchError):
        engine.render(template_path=str(path), context=RenderContext(data={}))


def test_missing_renderer_warn_policy_skips_placeholder(tmp_path: Path):
    template_path = tmp_path / "template.pptx"
    output_path = tmp_path / "output.pptx"
    logo_path = tmp_path / "logo.png"
    chart_path = tmp_path / "chart.png"
    _write_png(logo_path)
    _write_png(chart_path)
    _build_template(template_path)

    registry = RendererRegistry()
    registry.register_func("title", lambda placeholder, context: TextContent(text="已处理标题"))
    registry.register_func("logo", lambda placeholder, context: ImageContent(image_path=str(logo_path)))
    registry.register_func("sales", lambda placeholder, context: ChartContent(image_path=str(chart_path)))

    engine = PptTemplateEngine(registry, EngineOptions(missing_renderer_policy="warn"))
    result = engine.render(
        template_path=str(template_path),
        output_path=str(output_path),
        context=RenderContext(
            data={"project": {"name": "北极星"}, "report_date": "2026-03-31", "owner": {"name": "Alice"}}
        ),
    )

    assert result.success is True
    assert result.rendered_count == 3
    assert result.skipped_count == 1
    assert any("missing renderer for placeholder key 'risk_table'" in warning for warning in result.warnings)

    rendered = Presentation(str(output_path))
    slide = rendered.slides[0]
    title_shape = next(shape for shape in slide.shapes if getattr(shape, "name", None) == "ph:text:title")
    skipped_table_shape = next(shape for shape in slide.shapes if getattr(shape, "name", None) == "ph:table:risk_table")
    assert title_shape.text == "已处理标题"
    assert skipped_table_shape.text == "table"


def test_text_replacer_public_api(tmp_path: Path):
    template_path = tmp_path / "template.pptx"
    _build_template(template_path)
    prs = Presentation(str(template_path))
    replacer = TextReplacer()
    result = replacer.replace_presentation_text(
        prs,
        context=RenderContext(data={"project": {"name": "Aurora"}, "report_date": "2026-04-01", "owner": {"name": "Bob"}}),
    )

    texts = [shape.text for shape in prs.slides[0].shapes if getattr(shape, "has_text_frame", False)]
    assert "项目：Aurora 日期：2026-04-01 缺失：" in texts
    assert result.replaced_count >= 4
    assert any("missing text field 'missing.value'" in warning for warning in result.warnings)


def test_operations_slide_and_section_flow(tmp_path: Path):
    prs = Presentation()
    for idx in range(3):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1)).text = f"slide-{idx}"
    path = tmp_path / "ops.pptx"
    prs.save(path)

    ops = PptOperations.load(template_path=str(path))
    ops.add_section("Content", 1)
    ops.insert_slide(target_index=1, layout_index=6)
    ops.delete_slide(3)
    ops.delete_section(0)

    output = tmp_path / "ops-out.pptx"
    ops.save_to_path(str(output))
    rendered = Presentation(str(output))
    assert len(rendered.slides) == 3
    slide0_texts = [shape.text for shape in rendered.slides[0].shapes if getattr(shape, "has_text_frame", False)]
    slide1_texts = [shape.text for shape in rendered.slides[1].shapes if getattr(shape, "has_text_frame", False)]
    assert slide0_texts == ["slide-0"]
    assert slide1_texts == []

    with ZipFile(output) as zf:
        presentation_xml = zf.read("ppt/presentation.xml").decode("utf-8")
    assert 'name="Content"' in presentation_xml
    assert "sectionLst" in presentation_xml
    assert presentation_xml.count("<p:sldId ") == 3


def test_operations_table_row_column_and_merge(tmp_path: Path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = slide.shapes.add_table(3, 3, Inches(1), Inches(1), Inches(4), Inches(2))
    shape.name = "ops-table"
    table = shape.table
    for row_index in range(3):
        for col_index in range(3):
            table.cell(row_index, col_index).text = f"{row_index},{col_index}"
    path = tmp_path / "table-ops.pptx"
    prs.save(path)

    ops = PptOperations.load(template_path=str(path))
    ops.delete_table_row(0, "ops-table", 1)
    ops.delete_table_column(0, "ops-table", 1)
    ops.merge_table_cells(0, "ops-table", 0, 0, 1, 1)

    output_bytes = ops.save_to_bytes()
    rendered = Presentation(BytesIO(output_bytes))
    shape = rendered.slides[0].shapes[0]
    table = shape.table
    assert len(table.rows) == 2
    assert len(table.columns) == 2
    assert table.cell(0, 0).is_merge_origin is True


def test_row_column_delete_rejects_merged_tables(tmp_path: Path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = slide.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(4), Inches(2))
    shape.name = "merged-table"
    shape.table.cell(0, 0).merge(shape.table.cell(1, 1))
    path = tmp_path / "merged.pptx"
    prs.save(path)

    ops = PptOperations.load(template_path=str(path))
    with pytest.raises(OperationError):
        ops.delete_table_row(0, "merged-table", 0)


def test_singlefile_render_text_and_operations(tmp_path: Path):
    sdk = _load_singlefile_module()
    template_path = tmp_path / "template.pptx"
    output_path = tmp_path / "output.pptx"
    logo_path = tmp_path / "logo.png"
    chart_path = tmp_path / "chart.png"
    _write_png(logo_path)
    _write_png(chart_path)
    _build_template(template_path)

    registry = sdk.RendererRegistry()
    registry.register_func("title", lambda placeholder, context: sdk.TextContent(text="Single File"))
    registry.register_func("logo", lambda placeholder, context: sdk.ImageContent(image_path=str(logo_path)))
    registry.register_func("sales", lambda placeholder, context: sdk.ChartContent(image_path=str(chart_path)))
    registry.register_func("risk_table", lambda placeholder, context: sdk.TableContent(headers=["A"], rows=[["B"]]))

    engine = sdk.PptTemplateEngine(registry, sdk.EngineOptions())
    result = engine.render(
        template_path=str(template_path),
        output_path=str(output_path),
        context=sdk.RenderContext(data={"project": {"name": "Aurora"}, "report_date": "2026-04-01", "owner": {"name": "Bob"}}),
    )
    assert result.success is True

    text_template = tmp_path / "text-template.pptx"
    _build_template(text_template)
    prs = Presentation(str(text_template))
    replace_result = sdk.TextReplacer().replace_presentation_text(
        prs,
        context=sdk.RenderContext(data={"project": {"name": "Aurora"}, "owner": {"name": "Bob"}, "report_date": "2026-04-01"}),
    )
    assert replace_result.replaced_count >= 2

    ops = sdk.PptOperations.load(template_path=str(output_path))
    ops.insert_slide(target_index=1, layout_index=6)
    ops.add_section(name="正文", start_slide_index=1)
    final_bytes = ops.save_to_bytes()
    assert isinstance(final_bytes, bytes)


def test_singlefile_exports_match_expected_surface():
    sdk = _load_singlefile_module()
    expected = {
        "BaseRenderer",
        "ChartContent",
        "Content",
        "ContentTypeMismatchError",
        "cell",
        "DuplicatePlaceholderError",
        "EngineOptions",
        "FieldReplaceError",
        "ImageContent",
        "OperationError",
        "Placeholder",
        "PlaceholderFormatError",
        "PptTemplateEngine",
        "PptOperations",
        "PptTemplateSdkError",
        "RenderContext",
        "RenderResult",
        "RendererNotFoundError",
        "RendererRegistry",
        "ShapeOperationError",
        "TableCellsContent",
        "TableContent",
        "TemplateParseError",
        "TextReplaceResult",
        "TextReplacer",
        "TextContent",
        "ValidationReport",
    }
    assert set(sdk.__all__) == expected


def test_singlefile_register_func_supports_bound_kwargs(tmp_path: Path):
    sdk = _load_singlefile_module()
    template_path = tmp_path / "bound.pptx"
    _build_dual_text_template(template_path)

    registry = sdk.RendererRegistry()

    def render_with_prefix(placeholder, context, prefix):
        return sdk.TextContent(text=f"{prefix}{context.get_value('project.name')}")

    registry.register_func("title", render_with_prefix, prefix="标题：")
    registry.register_func("subtitle", render_with_prefix, prefix="副标题：")

    result = sdk.PptTemplateEngine(registry).render(
        template_path=str(template_path),
        context=sdk.RenderContext(data={"project": {"name": "北极星"}}),
    )

    rendered = Presentation(BytesIO(result.output_bytes))
    texts = [shape.text for shape in rendered.slides[0].shapes if getattr(shape, "has_text_frame", False)]
    assert "标题：北极星" in texts
    assert "副标题：北极星" in texts


def test_singlefile_text_placeholder_preserves_primary_style(tmp_path: Path):
    sdk = _load_singlefile_module()
    template_path = tmp_path / "styled-text.pptx"
    output_path = tmp_path / "styled-text-out.pptx"
    _build_styled_text_template(template_path)

    registry = sdk.RendererRegistry()
    registry.register_func("title", lambda placeholder, context: sdk.TextContent(text="新的标题"))

    sdk.PptTemplateEngine(registry).render(
        template_path=str(template_path),
        output_path=str(output_path),
        context=sdk.RenderContext(data={}),
    )

    rendered = Presentation(str(output_path))
    shape = next(shape for shape in rendered.slides[0].shapes if getattr(shape, "name", None) == "ph:text:title")
    paragraph = shape.text_frame.paragraphs[0]
    run = paragraph.runs[0]

    assert shape.text == "新的标题"
    assert paragraph.alignment == PP_ALIGN.CENTER
    assert run.font.name == "Arial"
    assert run.font.size == Pt(24)
    assert run.font.bold is True
    assert run.font.italic is True
    assert run.font.color.rgb == RGBColor(0x11, 0x22, 0x33)


def test_singlefile_native_table_placeholder_preserves_layout_and_style(tmp_path: Path):
    sdk = _load_singlefile_module()
    template_path = tmp_path / "native-table.pptx"
    output_path = tmp_path / "native-table-out.pptx"
    _build_native_table_placeholder_template(template_path)

    registry = sdk.RendererRegistry()
    registry.register_func(
        "risk_table",
        lambda placeholder, context: sdk.TableContent(headers=["风险", "等级"], rows=[["现金流", "高"]]),
    )

    sdk.PptTemplateEngine(registry).render(
        template_path=str(template_path),
        output_path=str(output_path),
        context=sdk.RenderContext(data={}),
    )

    rendered = Presentation(str(output_path))
    shape = next(shape for shape in rendered.slides[0].shapes if getattr(shape, "name", None) == "ph:table:risk_table")
    table = shape.table
    paragraph = table.cell(0, 0).text_frame.paragraphs[0]
    run = paragraph.runs[0]

    assert len(table.rows) == 2
    assert len(table.columns) == 2
    assert table.columns[0].width == Inches(3)
    assert table.columns[1].width == Inches(2)
    assert table.rows[0].height == Inches(0.8)
    assert table.rows[1].height == Inches(1.2)
    assert table.cell(1, 0).text == "现金流"
    assert paragraph.alignment == PP_ALIGN.CENTER
    assert run.font.name == "Arial"
    assert run.font.size == Pt(20)
    assert run.font.bold is True
    assert run.font.color.rgb == RGBColor(0x44, 0x55, 0x66)


def test_singlefile_native_table_placeholder_size_mismatch_errors(tmp_path: Path):
    sdk = _load_singlefile_module()
    template_path = tmp_path / "native-table-mismatch.pptx"
    _build_native_table_placeholder_template(template_path)

    registry = sdk.RendererRegistry()
    registry.register_func(
        "risk_table",
        lambda placeholder, context: sdk.TableContent(headers=["风险", "等级", "状态"], rows=[["现金流", "高", "开启"]]),
    )

    engine = sdk.PptTemplateEngine(registry)
    with pytest.raises(sdk.ShapeOperationError, match="table placeholder size mismatch"):
        engine.render(template_path=str(template_path), context=sdk.RenderContext(data={}))


def test_singlefile_table_cells_content_updates_only_targeted_cells(tmp_path: Path):
    sdk = _load_singlefile_module()
    template_path = tmp_path / "native-table-cells.pptx"
    output_path = tmp_path / "native-table-cells-out.pptx"
    _build_native_table_placeholder_template(template_path)

    registry = sdk.RendererRegistry()
    registry.register_func(
        "risk_table",
        lambda placeholder, context: sdk.TableCellsContent(cells={(1, 0): "现金流", (1, 1): "高"}),
    )

    sdk.PptTemplateEngine(registry).render(
        template_path=str(template_path),
        output_path=str(output_path),
        context=sdk.RenderContext(data={}),
    )

    rendered = Presentation(str(output_path))
    shape = next(shape for shape in rendered.slides[0].shapes if getattr(shape, "name", None) == "ph:table:risk_table")
    table = shape.table
    header_paragraph = table.cell(0, 0).text_frame.paragraphs[0]
    header_run = header_paragraph.runs[0]

    assert table.cell(0, 0).text == "风险"
    assert table.cell(0, 1).text == "等级"
    assert table.cell(1, 0).text == "现金流"
    assert table.cell(1, 1).text == "高"
    assert header_paragraph.alignment == PP_ALIGN.CENTER
    assert header_run.font.name == "Arial"
    assert header_run.font.size == Pt(20)
    assert header_run.font.bold is True
    assert header_run.font.color.rgb == RGBColor(0x44, 0x55, 0x66)


def test_singlefile_table_cells_content_rejects_non_native_table_placeholder(tmp_path: Path):
    sdk = _load_singlefile_module()
    template_path = tmp_path / "textbox-table-placeholder.pptx"
    _build_template(template_path)

    registry = sdk.RendererRegistry()
    registry.register_func("title", lambda placeholder, context: sdk.TextContent(text="标题"))
    registry.register_func("logo", lambda placeholder, context: sdk.ImageContent(image_path=str(tmp_path / "logo.png")))
    registry.register_func("sales", lambda placeholder, context: sdk.ChartContent(image_path=str(tmp_path / "chart.png")))
    registry.register_func("risk_table", lambda placeholder, context: sdk.TableCellsContent(cells={(0, 0): "仅更新"}))
    _write_png(tmp_path / "logo.png")
    _write_png(tmp_path / "chart.png")

    with pytest.raises(sdk.ShapeOperationError, match="cannot accept partial table cell updates"):
        sdk.PptTemplateEngine(registry).render(
            template_path=str(template_path),
            context=sdk.RenderContext(data={"project": {"name": "Aurora"}, "owner": {"name": "Bob"}, "report_date": "2026-04-01"}),
        )


def test_singlefile_patch_table_cells_operation_updates_selected_cells(tmp_path: Path):
    sdk = _load_singlefile_module()
    template_path = tmp_path / "ops-patch-table.pptx"
    _build_native_table_placeholder_template(template_path)

    ops = sdk.PptOperations.load(template_path=str(template_path))
    ops.patch_table_cells(0, "ph:table:risk_table", {(1, 0): "履约", (1, 1): ""})

    rendered = Presentation(BytesIO(ops.save_to_bytes()))
    shape = next(shape for shape in rendered.slides[0].shapes if getattr(shape, "name", None) == "ph:table:risk_table")
    table = shape.table

    assert table.cell(0, 0).text == "风险"
    assert table.cell(0, 1).text == "等级"
    assert table.cell(1, 0).text == "履约"
    assert table.cell(1, 1).text == ""


def test_singlefile_patch_table_cells_rejects_out_of_range_coordinates(tmp_path: Path):
    sdk = _load_singlefile_module()
    template_path = tmp_path / "ops-patch-table-range.pptx"
    _build_native_table_placeholder_template(template_path)

    ops = sdk.PptOperations.load(template_path=str(template_path))
    with pytest.raises(sdk.ShapeOperationError, match="out of range"):
        ops.patch_table_cells(0, "ph:table:risk_table", {(9, 9): "bad"})


def test_table_content_can_override_cell_font_style(tmp_path: Path):
    template_path = tmp_path / "styled-table.pptx"
    output_path = tmp_path / "styled-table-out.pptx"
    _build_native_table_placeholder_template(template_path)

    registry = RendererRegistry()
    registry.register_func(
        "risk_table",
        lambda placeholder, context: TableContent(
            headers=["风险", "等级"],
            rows=[
                [
                    "现金流",
                    cell("高", color="FF0000", bold=True),
                ]
            ],
        ),
    )

    PptTemplateEngine(registry).render(
        template_path=str(template_path),
        output_path=str(output_path),
        context=RenderContext(data={}),
    )

    rendered = Presentation(str(output_path))
    shape = next(shape for shape in rendered.slides[0].shapes if getattr(shape, "name", None) == "ph:table:risk_table")
    target_run = shape.table.cell(1, 1).text_frame.paragraphs[0].runs[0]

    assert shape.table.cell(1, 1).text == "高"
    assert target_run.font.bold is True
    assert target_run.font.color.rgb == RGBColor(0xFF, 0x00, 0x00)
    assert target_run.font.name is None
    assert target_run.font.size == Pt(12)


def test_table_cells_content_can_override_partial_cell_font_style(tmp_path: Path):
    template_path = tmp_path / "partial-styled-table.pptx"
    output_path = tmp_path / "partial-styled-table-out.pptx"
    _build_native_table_placeholder_template(template_path)

    registry = RendererRegistry()
    registry.register_func(
        "risk_table",
        lambda placeholder, context: TableCellsContent(
            cells={
                (1, 1): cell("中", color="00AA00", italic=True)
            }
        ),
    )

    PptTemplateEngine(registry).render(
        template_path=str(template_path),
        output_path=str(output_path),
        context=RenderContext(data={}),
    )

    rendered = Presentation(str(output_path))
    shape = next(shape for shape in rendered.slides[0].shapes if getattr(shape, "name", None) == "ph:table:risk_table")
    target_run = shape.table.cell(1, 1).text_frame.paragraphs[0].runs[0]

    assert shape.table.cell(1, 0).text == "旧值"
    assert shape.table.cell(1, 1).text == "中"
    assert target_run.font.italic is True
    assert target_run.font.color.rgb == RGBColor(0x00, 0xAA, 0x00)
    assert target_run.font.size == Pt(12)


def test_table_cells_content_append_adds_new_run(tmp_path: Path):
    template_path = tmp_path / "append-table.pptx"
    output_path = tmp_path / "append-table-out.pptx"
    _build_native_table_placeholder_template(template_path)

    registry = RendererRegistry()
    registry.register_func(
        "risk_table",
        lambda placeholder, context: TableCellsContent(
            cells={
                (1, 1): cell(" +高", color="FF0000", bold=True, append=True)
            }
        ),
    )

    PptTemplateEngine(registry).render(
        template_path=str(template_path),
        output_path=str(output_path),
        context=RenderContext(data={}),
    )

    rendered = Presentation(str(output_path))
    shape = next(shape for shape in rendered.slides[0].shapes if getattr(shape, "name", None) == "ph:table:risk_table")
    paragraph = shape.table.cell(1, 1).text_frame.paragraphs[0]
    appended_run = paragraph.runs[-1]

    assert shape.table.cell(1, 1).text == "旧值 +高"
    assert len(paragraph.runs) >= 2
    assert appended_run.text == " +高"
    assert appended_run.font.bold is True
    assert appended_run.font.color.rgb == RGBColor(0xFF, 0x00, 0x00)
    assert appended_run.font.size == Pt(12)
